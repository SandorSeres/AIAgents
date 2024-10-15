from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
"""
A JSON struktúra legyen egy lista, amely minden egyes diát reprezentál egy objektumként. Minden diához tartozzon:

type: A diatípus, például "title" vagy "content".
title: A dia címe.
subtitle: (Opcionális) Csak a cím diákhoz szükséges, a felirat vagy további információk.
name: (Opcionális) Az előadó neve, ha szükséges.
date: (Opcionális) A dátum, ha szükséges.
bullet_points: (Opcionális) Lista a felsorolási pontokból.
image_path: (Opcionális) Kép elérési útja, ha szükséges.
Példa JSON fájl:
"""
# A PPTX fájl elérési útja
pptx_file = "./Generativ_AI_HR_Egeszsegugy.pptx"

# Ellenőrizzük, hogy a PPTX fájl létezik-e
if os.path.exists(pptx_file):
    # Ha létezik, betöltjük a meglévő prezentációt
    prs = Presentation(pptx_file)
else:
    # Ha nem létezik, létrehozunk egy új prezentációt
    prs = Presentation()


# Define the slide layout (Title Slide, Text Slide)
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Define a function to set black background and white text
def set_black_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

def add_title_slide(title, subtitle, name, date):
    slide = prs.slides.add_slide(title_slide_layout)
    set_black_background(slide)
    
    title_shape = slide.shapes.title
    title_shape.text = title

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = f"{subtitle}\n\nElőadó: {name}\nDátum: {date}"
    
    # Set font color to white
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

def add_content_slide(title, bullet_points, image_path=None):
    slide = prs.slides.add_slide(content_slide_layout)
    set_black_background(slide)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Ensure the title text is set to white
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)

    text_shape = slide.shapes.placeholders[1].text_frame
    for point in bullet_points:
        p = text_shape.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
    
    # Set font color to white for bullet points
    for paragraph in text_shape.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)


# Add title slide
add_title_slide(
    "Generatív Mesterséges Intelligencia Alkalmazása a HR Területén az Egészségügyben", 
    "A Toborzás, Képzés és Hatékonyság Forradalmasítása", 
    "[Az Ön neve]", 
    "[Dátum]"
)

# Slide 2: Bevezetés
add_content_slide(
    "Bevezetés",
    [
        "A generatív AI forradalmasítja a HR menedzsmentet az egészségügyben.",
        "Főbb hatásterületek: toborzás, képzés, munkaerő-hatékonyság, betegellátás."
    ]
)

# Slide 3: Miért Szükséges AI a HR-ben az Egészségügyben?
add_content_slide(
    "Miért Szükséges AI a HR-ben az Egészségügyben?",
    [
        "Az egészségügyi munkaerő komplexitása: orvosok, ápolók, adminisztratív személyzet kezelése.",
        "Hatékonyság iránti igény, jobb betegellátás és költségcsökkentés.",
        "Az AI automatizált megoldásokat kínál a HR folyamatok egyszerűsítésére."
    ]
)

# Slide 4: Toborzási Folyamatok Automatizálása
add_content_slide(
    "Toborzási Folyamatok Automatizálása",
    [
        "A generatív AI automatizálja a jelöltek keresését és az önéletrajzok szűrését.",
        "Adatelemzés a megfelelő jelöltek kiválasztására.",
        "Csökkenti a manuális szűrés idejét, a HR a legjobb jelöltekre fókuszálhat."
    ]
)

# Slide 5: Képzési Programok Személyre Szabása
add_content_slide(
    "Képzési Programok Személyre Szabása",
    [
        "AI-vezérelt tanulási utak, amelyek személyre szabottak az alkalmazottak számára.",
        "Folyamatos teljesítményértékelés és adaptív tartalomajánlás.",
        "Segíti az alkalmazottak fejlődését valós idejű adatok alapján."
    ]
)

# Slide 6: Munkaerő Hatékonyságának Növelése
add_content_slide(
    "Munkaerő Hatékonyságának Növelése",
    [
        "A generatív AI automatizálja a rutinfeladatokat (pl. bérszámfejtés, beillesztési folyamatok).",
        "Több idő marad a stratégiai HR tervezésre és fontosabb feladatokra.",
        "Csökkenti a munkaerő-terheket és növeli a hatékonyságot."
    ]
)

# Slide 7: Betegellátás és Adminisztráció
add_content_slide(
    "Betegellátás és Adminisztráció",
    [
        "Az AI automatizálja az adminisztratív feladatokat, például a betegadatok kezelését.",
        "Javítja a munkafolyamatokat, enyhítve az orvosi személyzet terheit.",
        "Optimalizálja a kezelési terveket, időt és költségeket takarít meg."
    ]
)

# Adding slide 8 to 15 to the presentation

# Slide 8: Képzés és Szimuláció
add_content_slide(
    "Képzés és Szimuláció",
    [
        "A generatív AI szimulációs képzéseket biztosít kockázatmentes környezetben.",
        "Az orvostanhallgatók biztonságosan gyakorolhatják az eljárásokat.",
        "Javítja a készségeket és felkészíti a dolgozókat a valós helyzetekre."
    ]
)

# Slide 9: Döntéstámogatás a HR-ben
add_content_slide(
    "Döntéstámogatás a HR-ben",
    [
        "Az AI alapú döntéstámogatási rendszerek segítik a HR szakembereket és az orvosokat.",
        "Jobb döntések a személyzet kezelésében, javítva a betegellátás minőségét."
    ]
)

# Slide 10: Költségek Csökkentése
add_content_slide(
    "Költségek Csökkentése",
    [
        "Az AI előrejelzi a berendezések karbantartását, elkerülve a váratlan leállásokat.",
        "Költségcsökkentés a hatékonyság és a munkaerőkezelés optimalizálásával."
    ]
)

# Slide 11: Betegélmény Javítása
add_content_slide(
    "Betegélmény Javítása",
    [
        "Az AI által vezérelt chatbotok segítenek az időpontfoglalásban és egyéb adminisztrációs feladatokban.",
        "Javítják a betegek tapasztalatait azáltal, hogy gyorsabb és egyszerűbb szolgáltatásokat biztosítanak."
    ]
)

# Slide 12: Kockázatkezelés
add_content_slide(
    "Kockázatkezelés",
    [
        "Az AI segít az egészségügyi kockázatok azonosításában és kezelésében.",
        "Lehetőséget ad proaktív intézkedések megtételére a betegbiztonság növelése érdekében."
    ]
)

# Slide 13: Jövőbeli Kilátások
add_content_slide(
    "Jövőbeli Kilátások",
    [
        "A generatív AI jövője az egészségügyben: korai betegségmegelőzés, személyre szabott kezelési tervek.",
        "Az AI tovább javítja a betegellátás minőségét és az egészségügyi folyamatokat."
    ]
)

# Slide 14: Esettanulmányok és Példák
add_content_slide(
    "Esettanulmányok és Példák",
    [
        "Rövid bemutatás konkrét esetekről, ahol az AI sikeresen alkalmazható a HR-ben az egészségügyben.",
        "Példák a toborzás, képzés, és a munkaerő-hatékonyság javítására."
    ]
)

# Slide 15: Következtetés
add_content_slide(
    "Következtetés",
    [
        "A generatív AI forradalmasítja a HR menedzsmentet és az egészségügyi szektort.",
        "Hatékonyabb munkafolyamatok, jobb betegellátás, költségcsökkentés és javított betegélmények.",
        "Az AI alkalmazása nélkülözhetetlen a jövő egészségügyi HR menedzsmentjében."
    ]
)

# Save the updated presentation with all slides
pptx_file = "./Generativ_AI_HR_Egeszsegugy.pptx"
prs.save(pptx_file)
