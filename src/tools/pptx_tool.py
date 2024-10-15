from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

# Helper function to ensure directory exists
def check_path(path):
    if "/path/to/directory" in path:
        return False
    return True

# Define a function to set black background and white text
def set_black_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)


class ReadPPTXTool:
    """
    Class Name: ReadPPTXTool
    Description: A utility class designed to read the content of a specified PowerPoint (.pptx) file.
    
    Attributes:
        name (str): The name of the tool.
        description (str): A brief description of what the tool does.
        parameters (list): A list of parameters required by the tool, including the filename and directory.
    
    Methods:
        _run(filename, directory): Reads the content of the specified PowerPoint file and returns the text or an error message.
    """
    name: str = "ReadPPTXTool"
    description: str = "A tool to read the content of a PowerPoint (.pptx) file."
    parameters: list = ["filename", "directory"]
    parameters: str = "Mandatory: filename (str), directory (str)"

    def _run(self, filename: str, directory: str) -> tuple:
        try:
            # Construct file path
            file_path = os.path.join(directory, filename)
            # Ensure the file exists
            if not os.path.exists(file_path):
                return f"File {filename} does not exist in directory {directory}.", False

            # Load the presentation
            prs = Presentation(file_path)
            slides_content = []

            # Extract text from slides
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slides_content.append("\n".join(slide_text))

            return "\n---Slide---\n".join(slides_content), True
        except Exception as e:
            return f"An error occurred: {str(e)}", False
    def clone(self):
        """
        Creates a clone of the ReadPPTXTool instance.

        Returns:
            ReadPPTXTool: A new instance of ReadPPTXTool.
        """
        return ReadPPTXTool()


class SaveToPPTXTool:
    """
    Class Name: SaveToPPTXTool
    Description: A utility class to save content to a new PowerPoint (.pptx) file.

    Methods:
        _run(content, filename, directory): Saves the provided content as a new PowerPoint presentation.
    """
    name: str = "SaveToPPTXTool"
    description: str = "A tool to save content as a new PowerPoint presentation."
    parameters: str = """Mandatory: content (dict) , filename (str), directory (str)
    Example content:
        content = [
        {
            "type": "title",
            "title": "Generative AI in Healthcare HR",
        },
        {
            "type": "content",
            "title": "Introduction to AI in HR",
            "bullet_points": [
                "AI is revolutionizing HR processes.",
                "Automating recruitment, training, and efficiency improvement."
            ]
        },
    ]
    """

    def _run(self, content: dict, filename: str, directory: str) -> tuple:
        try:
            if not check_path(directory):
                return "\'/path/to/directory\' is not a valid directory. Please try again!", False

            # Ensure the directory exists
            os.makedirs(directory, exist_ok=True)

            # Define the file path
            file_path = os.path.join(directory, filename)
            prs = Presentation()

            # Add slides based on the content structure
            for slide_data in content:
                slide_type = slide_data.get("type", "content")
                title = slide_data.get("title", "")
                bullet_points = slide_data.get("bullet_points", [])
                
                if slide_type == "title":
                    slide_layout = prs.slide_layouts[0]
                    set_black_background(slide)
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = title
                else:
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    set_black_background(slide)
                    slide.shapes.title.text = title
                    text_frame = slide.shapes.placeholders[1].text_frame
                    for point in bullet_points:
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.font.size = Pt(24)

            # Save the presentation
            prs.save(file_path)
            return f"Presentation saved as {file_path}", True

        except Exception as e:
            return f"An error occurred: {str(e)}", False
    def clone(self):
        """
        Creates a clone of the SaveToPPTXTool instance.

        Returns:
            SaveToPPTXTool: A new instance of SaveToPPTXTool.
        """
        return SaveToPPTXTool()


class AppendToPPTXTool:
    """
    Class Name: AppendToPPTXTool
    Description: A utility class to append content to an existing PowerPoint (.pptx) file.

    Methods:
        _run(content, filename, directory): Appends the provided content to an existing PowerPoint presentation.
    """
    name: str = "AppendToPPTXTool"
    description: str = "A tool to append content to an existing PowerPoint presentation."
    parameters: str = """Mandatory: content (dict) , filename (str), directory (str)
    Example content:
        content = [
        {
            "type": "title",
            "title": "Generative AI in Healthcare HR",
        },
        {
            "type": "content",
            "title": "Introduction to AI in HR",
            "bullet_points": [
                "AI is revolutionizing HR processes.",
                "Automating recruitment, training, and efficiency improvement."
            ]
        },
    ]
    """

    def _run(self, content: dict, filename: str, directory: str) -> tuple:
        try:
            if not check_path(directory):
                return "\'/path/to/directory\' is not a valid directory. Please try again!", False

            # Ensure the directory exists or create it
            os.makedirs(directory, exist_ok=True)

            # Define the file path
            file_path = os.path.join(directory, filename)

            # Check if the file exists, if not, create a new empty presentation
            if not os.path.exists(file_path):
                prs = Presentation()  # Create a new empty presentation
            else:
                prs = Presentation(file_path)  # Load the existing presentation

            # Append new slides based on the content
            for slide_data in content:
                slide_type = slide_data.get("type", "content")
                title = slide_data.get("title", "")
                bullet_points = slide_data.get("bullet_points", [])
                
                if slide_type == "title":
                    slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(slide_layout)
                    set_black_background(slide)
                    slide.shapes.title.text = title
                else:
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    set_black_background(slide)
                    slide.shapes.title.text = title
                    text_frame = slide.shapes.placeholders[1].text_frame
                    for point in bullet_points:
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.font.size = Pt(24)

            # Save the updated presentation
            prs.save(file_path)
            return f"Appended to presentation: {file_path}", True

        except Exception as e:
            return f"An error occurred: {str(e)}", False
    def clone(self):
        """
        Creates a clone of the AppendToPPTXTool instance.

        Returns:
            AppendToPPTXTool: A new instance of AppendToPPTXTool.
        """
        return AppendToPPTXTool()
